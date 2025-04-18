#!/usr/bin/env python3

import argparse
import json
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def apply_style(style):
    """Apply style settings based on the selected theme"""
    styles = {
        "Business Professional": {
            "bg_color": RGBColor(255, 255, 255),
            "title_color": RGBColor(31, 73, 125),
            "text_color": RGBColor(0, 0, 0),
            "accent_color": RGBColor(79, 129, 189)
        },
        "Creative & Modern": {
            "bg_color": RGBColor(242, 242, 242),
            "title_color": RGBColor(255, 0, 110),
            "text_color": RGBColor(51, 51, 51),
            "accent_color": RGBColor(131, 56, 236)
        },
        "Academic": {
            "bg_color": RGBColor(255, 255, 255),
            "title_color": RGBColor(0, 51, 102),
            "text_color": RGBColor(0, 0, 0),
            "accent_color": RGBColor(0, 102, 204)
        },
        "Minimalist": {
            "bg_color": RGBColor(250, 250, 250),
            "title_color": RGBColor(80, 80, 80),
            "text_color": RGBColor(100, 100, 100),
            "accent_color": RGBColor(200, 200, 200)
        },
        "Bold & Vibrant": {
            "bg_color": RGBColor(0, 0, 0),
            "title_color": RGBColor(255, 190, 11),
            "text_color": RGBColor(255, 255, 255),
            "accent_color": RGBColor(251, 86, 7)
        },
    }
    
    return styles.get(style, styles["Business Professional"])

def create_presentation(title, slides_data, style, output_path):
    try:
        # Create a presentation
        prs = Presentation()
        
        # Get style settings
        theme = apply_style(style)
        
        # Create title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Set the background color
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = theme["bg_color"]
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme["title_color"]
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Add subtitle if needed
        if 1 in slide.placeholders:
            subtitle = slide.placeholders[1]
            subtitle.text = f"Style: {style}"
            subtitle.text_frame.paragraphs[0].font.color.rgb = theme["accent_color"]
            subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        
        # Create content slides
        for slide_data in slides_data:
            content_slide_layout = prs.slide_layouts[1]  # Using layout with title and content
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Set background color
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = theme["bg_color"]
            
            # Add slide title
            title_shape = slide.shapes.title
            title_shape.text = slide_data.get("title", "")
            title_shape.text_frame.paragraphs[0].font.color.rgb = theme["title_color"]
            title_shape.text_frame.paragraphs[0].font.size = Pt(36)
            title_shape.text_frame.paragraphs[0].font.bold = True
            
            # Add content
            content = slide_data.get("content", [])
            
            # Handle both list and string content types
            if 1 in slide.placeholders:
                body_shape = slide.placeholders[1]
                text_frame = body_shape.text_frame
                
                # Safely clear text frame if method exists
                if hasattr(text_frame, 'clear'):
                    text_frame.clear()
                
                if isinstance(content, list):
                    # Handle bullet points - with explicit bullet configuration
                    for i, point in enumerate(content):
                        p = text_frame.add_paragraph()
                        p.text = point
                        p.font.color.rgb = theme["text_color"]
                        p.font.size = Pt(24)
                        p.level = 0  # Set paragraph level
                        
                        # Explicitly enable bullet points - this is the key addition
                        p.bullet.visible = True
                else:
                    # Handle paragraph text
                    p = text_frame.add_paragraph()
                    p.text = content if isinstance(content, str) else str(content)
                    p.font.color.rgb = theme["text_color"]
                    p.font.size = Pt(24)
                    # No bullets for non-list content
            
            # Add notes if available
            if "notes" in slide_data:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = slide_data["notes"]
        
        # Save the presentation
        prs.save(output_path)
        print(f"Presentation saved to {output_path}")
        return output_path
        
    except Exception as e:
        print(f"Error in create_presentation: {str(e)}", file=sys.stderr)
        raise

def main():
    parser = argparse.ArgumentParser(description='Generate PowerPoint presentation from JSON data')
    parser.add_argument('--title', required=True, help='Presentation title')
    parser.add_argument('--style', required=True, help='Presentation style')
    parser.add_argument('--input', required=True, help='Input JSON file path')
    parser.add_argument('--output', required=True, help='Output PPTX file path')
    
    args = parser.parse_args()
    
    try:
        print(f"Loading slides data from {args.input}")
        # Load slides data from JSON file
        with open(args.input, 'r') as f:
            slides_data = json.load(f)
        
        print(f"Creating presentation with {len(slides_data)} slides")
        # Create the presentation
        create_presentation(args.title, slides_data, args.style, args.output)
        
    except Exception as e:
        print(f"Error generating presentation: {str(e)}", file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()